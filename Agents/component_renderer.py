"""
Component Renderer — composable visual building blocks for slides.

Each component renders itself into a rectangular region on a slide.
Native components use python-pptx shapes directly (editable).
Hybrid components will render SVG→PNG and embed as images in later phases.

Usage from pptx_renderer:
    from Agents.component_renderer import render_composition
    render_composition(slide, composition, content, design, deck_theme)
"""

import os
import copy
from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from Agents.icon_manager import resolve_icon_for_bullet


# ---------------------------------------------------------------------------
# Region definitions: (left, top, width, height) in inches
# ---------------------------------------------------------------------------

REGIONS: dict[str, tuple[float, float, float, float]] = {
    "full": (0.5, 1.6, 12.3, 5.4),
    "main": (0.5, 1.6, 7.8, 5.4),
    "side_note": (8.5, 1.6, 4.3, 5.4),
    "left": (0.5, 1.6, 5.9, 5.4),
    "right": (6.8, 1.6, 6.0, 5.4),
    "top_strip": (0.5, 1.5, 12.3, 1.6),
    "center": (2.0, 2.5, 9.3, 3.5),
    # Modern asymmetric regions
    "hero_left": (0.0, 0.0, 5.6, 7.5),
    "hero_right": (7.7, 0.0, 5.633, 7.5),
    "content_right": (6.1, 1.2, 6.7, 5.8),
    "content_left": (0.5, 1.2, 6.7, 5.8),
    "stat_main": (0.5, 1.0, 12.3, 3.5),
    "stat_context": (0.5, 5.0, 12.3, 2.0),
    "overlay_full": (0.0, 0.0, 13.333, 7.5),
}


# ---------------------------------------------------------------------------
# Layout engine — density-aware region computation  (Phase 4B)
# ---------------------------------------------------------------------------

# Slide canvas (inches) — below title bar
_CANVAS_LEFT = 0.5
_CANVAS_TOP = 1.6
_CANVAS_W = 12.3
_CANVAS_H = 5.4


def compute_adaptive_regions(
    components: list[dict],
    content: dict,
) -> dict[str, tuple[float, float, float, float]]:
    """Return a region map with sizes adapted to content density.

    Heuristics used:
    * Single full-region component → use entire canvas.
    * main + side_note → split at 65/35 default; widen main to 75% if
      side_note text is ≤2 bullets and main is a complex diagram.
    * top_strip + full → give the strip slightly more height when there are
      5+ KPI items and compress the body region below.
    * Bullet count >6 → expand height for text-heavy components.
    * If an image is present alongside a diagram, shrink diagram region and
      add an image pocket.
    """
    # Start from the defaults
    regions: dict[str, tuple[float, float, float, float]] = dict(REGIONS)

    if not components:
        return regions

    region_names = {c.get("region", "full") for c in components}
    comp_types = {c.get("type", "") for c in components}

    bullets = content.get("bullets", [])
    n_bullets = len(bullets)
    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})

    # --- Rule 1: main + side_note split ---
    if "main" in region_names and "side_note" in region_names:
        main_comp = next((c for c in components if c.get("region") == "main"), None)
        side_comp = next(
            (c for c in components if c.get("region") == "side_note"), None
        )
        main_type = main_comp.get("type", "") if main_comp else ""
        side_type = side_comp.get("type", "") if side_comp else ""

        complex_diagrams = {
            "timeline",
            "process_flow",
            "cycle_loop",
            "comparison_matrix",
            "value_chain",
            "card_grid",
            "chart_panel",
        }

        # Widen main for complex diagrams with light side text
        if main_type in complex_diagrams and side_type in ("text_block", "bullet_list"):
            side_bullets = n_bullets if side_type == "bullet_list" else 0
            if side_bullets <= 2:
                # 75/25 split
                main_w = _CANVAS_W * 0.73
                side_w = _CANVAS_W * 0.25
                gap = _CANVAS_W * 0.02
                regions["main"] = (_CANVAS_LEFT, _CANVAS_TOP, main_w, _CANVAS_H)
                side_left = _CANVAS_LEFT + main_w + gap
                regions["side_note"] = (side_left, _CANVAS_TOP, side_w, _CANVAS_H)
            else:
                # Standard 65/35 with slight adjustment for 3-4 bullets
                main_w = _CANVAS_W * 0.63
                side_w = _CANVAS_W * 0.35
                gap = _CANVAS_W * 0.02
                regions["main"] = (_CANVAS_LEFT, _CANVAS_TOP, main_w, _CANVAS_H)
                side_left = _CANVAS_LEFT + main_w + gap
                regions["side_note"] = (side_left, _CANVAS_TOP, side_w, _CANVAS_H)

    # --- Rule 2: top_strip + body ---
    if "top_strip" in region_names:
        kpi_comp = next((c for c in components if c.get("region") == "top_strip"), None)
        kpi_items = len(visual_data.get("kpi_metrics", []))
        if kpi_items >= 5:
            # Taller strip for many KPIs
            strip_h = 2.0
        elif kpi_items <= 2:
            strip_h = 1.2
        else:
            strip_h = 1.6
        regions["top_strip"] = (_CANVAS_LEFT, 1.5, _CANVAS_W, strip_h)
        # Push body content below strip
        body_top = 1.5 + strip_h + 0.15
        body_h = max(1.5, _CANVAS_TOP + _CANVAS_H - body_top)
        regions["full"] = (_CANVAS_LEFT, body_top, _CANVAS_W, body_h)
        regions["main"] = (_CANVAS_LEFT, body_top, _CANVAS_W * 0.63, body_h)

    # --- Rule 3: text-heavy full region → expand height by nudging top up ---
    if "full" in region_names and len(region_names) == 1:
        if n_bullets > 6:
            # Bring content region slightly higher / taller
            regions["full"] = (_CANVAS_LEFT, 1.4, _CANVAS_W, 5.7)

    # --- Rule 4: center region for hero_text → keep it vertically centered ---
    if "center" in region_names:
        title_len = len(content.get("title", ""))
        text_len = len(content.get("subtitle", content.get("key_takeaway", "")))
        if title_len + text_len > 200:
            # More vertical room for long text
            regions["center"] = (1.5, 2.0, 10.3, 4.5)
        elif title_len < 40:
            # Tighter for punchy short text
            regions["center"] = (2.5, 2.8, 8.3, 2.8)

    return regions


# ---------------------------------------------------------------------------
# Color helpers (duplicated from pptx_renderer to avoid circular imports)
# ---------------------------------------------------------------------------


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        hex_color = "2C3E50"
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _lighten(hex_color: str, factor: float = 0.85) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        hex_color = "2C3E50"
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return RGBColor(min(r, 255), min(g, 255), min(b, 255))


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=14,
    bold=False,
    color=None,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
    anchor=MSO_ANCHOR.TOP,
):
    """Add a text box with word wrap."""
    color = color or RGBColor(0x33, 0x33, 0x33)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


# ---------------------------------------------------------------------------
# Gradient fill helper
# ---------------------------------------------------------------------------


def _apply_gradient_fill(shape, hex_start: str, hex_end: str, angle: int = 0):
    """Apply a linear gradient fill to a shape via XML injection.

    angle: 0 = left→right, 5400000 = top→bottom (60000ths of a degree).
    """
    spPr = shape._element.spPr
    # Remove any existing fill
    for child in list(spPr):
        if (
            child.tag.endswith("}solidFill")
            or child.tag.endswith("}gradFill")
            or child.tag.endswith("}noFill")
        ):
            spPr.remove(child)

    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    gradFill = etree.SubElement(spPr, qn("a:gradFill"), attrib={"rotWithShape": "1"})
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

    # Stop 1 (start color at 0%)
    gs1 = etree.SubElement(gsLst, qn("a:gs"), attrib={"pos": "0"})
    srgb1 = etree.SubElement(
        gs1, qn("a:srgbClr"), attrib={"val": hex_start.lstrip("#")}
    )

    # Stop 2 (end color at 100%)
    gs2 = etree.SubElement(gsLst, qn("a:gs"), attrib={"pos": "100000"})
    srgb2 = etree.SubElement(gs2, qn("a:srgbClr"), attrib={"val": hex_end.lstrip("#")})

    lin = etree.SubElement(
        gradFill, qn("a:lin"), attrib={"ang": str(angle), "scaled": "1"}
    )


# ---------------------------------------------------------------------------
# Drop shadow helper
# ---------------------------------------------------------------------------


def _apply_shadow(shape, blur_pt: int = 6, offset_pt: int = 3, opacity_pct: int = 35):
    """Apply an outer drop shadow to a shape via XML injection."""
    spPr = shape._element.spPr
    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))

    outerShdw = etree.SubElement(
        effectLst,
        qn("a:outerShdw"),
        attrib={
            "blurRad": str(blur_pt * 12700),  # EMUs
            "dist": str(offset_pt * 12700),
            "dir": "5400000",  # bottom
            "rotWithShape": "0",
        },
    )
    srgb = etree.SubElement(outerShdw, qn("a:srgbClr"), attrib={"val": "000000"})
    etree.SubElement(srgb, qn("a:alpha"), attrib={"val": str(opacity_pct * 1000)})


# ---------------------------------------------------------------------------
# Decorative motifs
# ---------------------------------------------------------------------------


def _add_dot_pattern(
    slide,
    x_start: float,
    y_start: float,
    cols: int,
    rows: int,
    spacing: float,
    dot_size: float,
    hex_color: str,
    opacity: float = 0.3,
):
    """Add a subtle grid of small circles as a decorative background motif."""
    rgb = _hex_to_rgb(hex_color)
    for r in range(rows):
        for c in range(cols):
            cx = x_start + c * spacing
            cy = y_start + r * spacing
            dot = slide.shapes.add_shape(
                9, Inches(cx), Inches(cy), Inches(dot_size), Inches(dot_size)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = _lighten(hex_color, 1.0 - opacity)
            dot.line.fill.background()


def _add_decorative_circles(slide, accent: str, deck_theme: dict | None = None):
    """Add 2-3 overlapping translucent circles in a corner for visual interest."""
    secondary = _get_secondary({}, deck_theme) if deck_theme else accent
    configs = [
        (11.8, 5.8, 1.2, accent, 0.12),
        (12.2, 5.4, 0.9, secondary, 0.10),
        (11.4, 6.2, 0.6, accent, 0.08),
    ]
    for cx, cy, size, color, op in configs:
        circle = slide.shapes.add_shape(
            9, Inches(cx), Inches(cy), Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _lighten(color, 1.0 - op)
        circle.line.fill.background()


# Active adaptive region map, set by render_composition per slide
_active_region_map: dict | None = None


def _resolve_region(region_name: str, region_map: dict | None = None) -> tuple:
    """Convert region name to (left, top, width, height) in Inches.

    If *region_map* is provided (from compute_adaptive_regions) use it;
    otherwise check the module-level _active_region_map set by
    render_composition; finally fall back to the static REGIONS table.
    """
    source = region_map or _active_region_map or REGIONS
    coords = source.get(region_name, REGIONS.get(region_name, REGIONS["full"]))
    return (Inches(coords[0]), Inches(coords[1]), Inches(coords[2]), Inches(coords[3]))


def _get_accent(design: dict, deck_theme: dict | None = None) -> str:
    accent = design.get("color_accent", "#2C3E50")
    if deck_theme:
        accent = deck_theme.get("palette", {}).get("primary", accent)
    return accent


def _get_secondary(design: dict, deck_theme: dict | None = None) -> str:
    if deck_theme:
        return deck_theme.get("palette", {}).get("secondary", "#3498DB")
    return "#3498DB"


def _get_font(design: dict, element: str, deck_theme: dict | None = None) -> dict:
    """Get font spec from design hierarchy, with deck_theme fallback."""
    hierarchy = design.get("text_hierarchy", {})
    spec = hierarchy.get(element, {})
    defaults = {"title": 28, "subtitle": 16, "body": 14, "caption": 10, "metric": 36}
    heading_font = "Calibri"
    body_font = "Calibri"
    if deck_theme:
        fp = deck_theme.get("font_pair", {})
        heading_font = fp.get("heading", "Calibri")
        body_font = fp.get("body", "Calibri")
    is_heading = element in ("title", "metric")
    default_font = heading_font if is_heading else body_font
    return {
        "size": spec.get("size", defaults.get(element, 14)),
        "bold": spec.get("bold", element in ("title", "metric")),
        "font": spec.get("font", default_font),
        "color": spec.get("color", "#333333"),
    }


# ---------------------------------------------------------------------------
# Component renderers
# ---------------------------------------------------------------------------


def _render_hero_text(slide, region_name, content, design, props, deck_theme):
    """Large centered statement / thesis / quote."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    text = props.get("text") or content.get("subtitle", "") or ""
    if not text and content.get("bullets"):
        text = content["bullets"][0]

    _add_textbox(
        slide,
        left,
        top + Inches(0.5),
        width,
        height - Inches(1.0),
        text,
        font_size=32,
        bold=True,
        color=_hex_to_rgb(accent),
        alignment=PP_ALIGN.CENTER,
        font_name=_get_font(design, "title", deck_theme)["font"],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _render_bullet_list(slide, region_name, content, design, props, deck_theme):
    """Classic icon-bullet list."""
    left, top, width, height = _resolve_region(region_name)
    bullets = content.get("bullets", [])
    if not bullets:
        return

    max_bullets = props.get("max_bullets", len(bullets))
    bullets = bullets[:max_bullets]
    accent = _get_accent(design, deck_theme)
    body_spec = _get_font(design, "body", deck_theme)
    icon_hints = design.get("icon_suggestions", [])
    accent_hex = accent.lstrip("#")

    icon_size = Inches(0.26)
    icon_gap = Inches(0.10)
    text_left = Inches(left.inches + icon_size.inches + icon_gap.inches)
    text_width = Inches(width.inches - icon_size.inches - icon_gap.inches)
    row_height = Inches(0.44)

    for i, bullet in enumerate(bullets):
        y = Inches(top.inches + i * row_height.inches)
        icon_path = resolve_icon_for_bullet(icon_hints, i, accent_hex)
        if icon_path and os.path.isfile(icon_path):
            slide.shapes.add_picture(icon_path, left, y, icon_size, icon_size)
        else:
            _add_textbox(
                slide,
                left,
                y,
                icon_size,
                Inches(0.3),
                "🔹",
                font_size=12,
                alignment=PP_ALIGN.CENTER,
            )

        _add_textbox(
            slide,
            text_left,
            y,
            text_width,
            Inches(0.36),
            bullet,
            font_size=body_spec["size"],
            bold=False,
            color=_hex_to_rgb(body_spec["color"]),
            font_name=body_spec["font"],
        )


def _render_text_block(slide, region_name, content, design, props, deck_theme):
    """Paragraph prose block."""
    left, top, width, height = _resolve_region(region_name)
    max_lines = props.get("max_lines", 8)
    body_spec = _get_font(design, "body", deck_theme)

    # Use subtitle or join bullets as paragraph text
    text = props.get("text", "")
    if not text:
        bullets = content.get("bullets", [])[:max_lines]
        text = "\n".join(bullets) if bullets else content.get("subtitle", "")

    _add_textbox(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size=body_spec["size"],
        bold=False,
        color=_hex_to_rgb(body_spec["color"]),
        font_name=body_spec["font"],
    )


def _render_kpi_strip(slide, region_name, content, design, props, deck_theme):
    """Row of metric cards with value + label."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    metric_spec = _get_font(design, "metric", deck_theme)

    # Get metrics from visual_data or supporting_data
    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})
    card_items = visual_data.get("card_items", [])

    # Fallback to supporting_data strings
    if not card_items:
        supporting = content.get("supporting_data", [])
        card_items = [{"label": s, "value": ""} for s in supporting[:5]]

    if not card_items:
        return

    cards = card_items[:5]
    num = len(cards)
    card_w = min(2.8, (width.inches - 0.2 * (num - 1)) / num)
    card_h = min(height.inches, 1.5)
    gap = 0.2
    total = num * card_w + (num - 1) * gap
    start_x = left.inches + (width.inches - total) / 2

    light_bg = _lighten(accent, 0.88)
    accent_rgb = _hex_to_rgb(accent)
    metric_color = (
        _hex_to_rgb(metric_spec["color"]) if metric_spec.get("color") else accent_rgb
    )

    for i, card in enumerate(cards):
        cx = Inches(start_x + i * (card_w + gap))

        # Card background (rounded rectangle = shape type 5)
        shape = slide.shapes.add_shape(5, cx, top, Inches(card_w), Inches(card_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = light_bg
        shape.line.color.rgb = _lighten(accent, 0.6)
        shape.line.width = Pt(1)
        _apply_shadow(shape, blur_pt=5, offset_pt=2, opacity_pct=20)

        # Top accent stripe on card
        stripe = slide.shapes.add_shape(1, cx, top, Inches(card_w), Inches(0.06))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent_rgb
        stripe.line.fill.background()

        if isinstance(card, str):
            parts = card.strip().split(" ", 1)
            value_text, label_text = parts[0], (parts[1] if len(parts) > 1 else "")
        else:
            value_text = str(card.get("value", card.get("label", "")))
            label_text = str(card.get("label", "")) if card.get("value") else ""

        # Value
        _add_textbox(
            slide,
            cx + Inches(0.1),
            top + Inches(0.15),
            Inches(card_w - 0.2),
            Inches(0.6),
            value_text,
            font_size=min(metric_spec["size"], 38),
            bold=True,
            color=metric_color,
            alignment=PP_ALIGN.CENTER,
            font_name=metric_spec["font"],
        )
        # Label
        if label_text:
            _add_textbox(
                slide,
                cx + Inches(0.1),
                top + Inches(0.85),
                Inches(card_w - 0.2),
                Inches(0.5),
                label_text,
                font_size=11,
                bold=False,
                color=RGBColor(0x55, 0x55, 0x55),
                alignment=PP_ALIGN.CENTER,
            )


def _render_timeline(slide, region_name, content, design, props, deck_theme):
    """Horizontal timeline with event nodes, connecting lines, and labels."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    secondary = _get_secondary(design, deck_theme)

    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})
    events = visual_data.get("timeline_events", [])

    # Fallback: try to parse bullets as events
    if not events:
        for b in content.get("bullets", [])[:8]:
            events.append({"date": "", "label": b})

    if not events:
        return

    events = events[:8]
    num = len(events)
    accent_rgb = _hex_to_rgb(accent)
    secondary_rgb = _hex_to_rgb(secondary)

    # Draw horizontal axis line with gradient feel (thick base + thin top)
    axis_y = top.inches + height.inches * 0.45
    line_base = slide.shapes.add_shape(
        1, left, Inches(axis_y - 0.02), width, Inches(0.08)
    )
    line_base.fill.solid()
    line_base.fill.fore_color.rgb = _lighten(accent, 0.75)
    line_base.line.fill.background()

    line_top = slide.shapes.add_shape(1, left, Inches(axis_y), width, Inches(0.03))
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = accent_rgb
    line_top.line.fill.background()

    # Draw event nodes along the axis
    node_r = 0.20
    spacing = width.inches / (num + 1)

    for i, evt in enumerate(events):
        cx = left.inches + spacing * (i + 1)

        # Vertical connector line from node
        conn_y_start = axis_y - 0.5 if i % 2 == 0 else axis_y + node_r + 0.05
        conn_y_end = axis_y - node_r - 0.05 if i % 2 == 0 else axis_y + 0.5
        conn = slide.shapes.add_shape(
            1,
            Inches(cx - 0.015),
            Inches(min(conn_y_start, conn_y_end)),
            Inches(0.03),
            Inches(abs(conn_y_end - conn_y_start)),
        )
        conn.fill.solid()
        conn.fill.fore_color.rgb = _lighten(accent, 0.6)
        conn.line.fill.background()

        # Circle node with shadow
        circle = slide.shapes.add_shape(
            9,
            Inches(cx - node_r),
            Inches(axis_y - node_r),
            Inches(node_r * 2),
            Inches(node_r * 2),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_rgb if i % 2 == 0 else secondary_rgb
        circle.line.fill.background()
        _apply_shadow(circle, blur_pt=4, offset_pt=2, opacity_pct=25)

        # Date — alternate above/below for visual interest
        date_text = str(evt.get("date", evt.get("year", "")))
        if date_text:
            date_y = axis_y - 1.1 if i % 2 == 0 else axis_y + 0.6
            _add_textbox(
                slide,
                Inches(cx - 0.7),
                Inches(date_y),
                Inches(1.4),
                Inches(0.35),
                date_text,
                font_size=10,
                bold=True,
                color=accent_rgb,
                alignment=PP_ALIGN.CENTER,
            )

        # Label — opposite side from date
        label = str(evt.get("label", evt.get("event", evt.get("title", ""))))
        if label:
            label_y = axis_y + 0.6 if i % 2 == 0 else axis_y - 1.3
            _add_textbox(
                slide,
                Inches(cx - 0.9),
                Inches(label_y),
                Inches(1.8),
                Inches(0.8),
                label,
                font_size=9,
                bold=False,
                color=RGBColor(0x33, 0x33, 0x33),
                alignment=PP_ALIGN.CENTER,
            )


def _render_process_flow(slide, region_name, content, design, props, deck_theme):
    """Horizontal step boxes connected by arrows with shadows."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    secondary = _get_secondary(design, deck_theme)

    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})
    steps = visual_data.get("process_steps", [])

    if not steps:
        for i, b in enumerate(content.get("bullets", [])[:6]):
            steps.append({"label": b, "order": i + 1})

    if not steps:
        return

    steps = steps[:6]
    num = len(steps)
    accent_rgb = _hex_to_rgb(accent)
    secondary_rgb = _hex_to_rgb(secondary)

    box_w = min(1.8, (width.inches - 0.5 * (num - 1)) / num)
    box_h = 1.2
    arrow_w = 0.4
    total = num * box_w + (num - 1) * arrow_w
    start_x = left.inches + (width.inches - total) / 2
    box_y = top.inches + (height.inches - box_h) / 2

    for i, step in enumerate(steps):
        bx = start_x + i * (box_w + arrow_w)

        # Step box (rounded rectangle) with shadow
        shape = slide.shapes.add_shape(
            5, Inches(bx), Inches(box_y), Inches(box_w), Inches(box_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _lighten(accent, 0.85)
        shape.line.color.rgb = accent_rgb
        shape.line.width = Pt(1.5)
        _apply_shadow(shape, blur_pt=4, offset_pt=2, opacity_pct=18)

        # Colored number badge (small circle at top of box)
        badge_sz = 0.35
        badge = slide.shapes.add_shape(
            9,
            Inches(bx + (box_w - badge_sz) / 2),
            Inches(box_y - badge_sz / 3),
            Inches(badge_sz),
            Inches(badge_sz),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent_rgb if i % 2 == 0 else secondary_rgb
        badge.line.fill.background()
        _add_textbox(
            slide,
            Inches(bx + (box_w - badge_sz) / 2),
            Inches(box_y - badge_sz / 3 + 0.03),
            Inches(badge_sz),
            Inches(badge_sz - 0.06),
            str(step.get("order", i + 1)),
            font_size=13,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER,
        )

        # Step label
        label = str(step.get("label", step.get("title", "")))
        _add_textbox(
            slide,
            Inches(bx + 0.08),
            Inches(box_y + 0.35),
            Inches(box_w - 0.16),
            Inches(0.75),
            label,
            font_size=9,
            bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
            alignment=PP_ALIGN.CENTER,
        )

        # Arrow between steps
        if i < num - 1:
            ax = bx + box_w
            ay = box_y + box_h / 2
            _add_textbox(
                slide,
                Inches(ax),
                Inches(ay - 0.15),
                Inches(arrow_w),
                Inches(0.3),
                "▶",
                font_size=16,
                bold=True,
                color=accent_rgb,
                alignment=PP_ALIGN.CENTER,
            )


def _render_cycle_loop(slide, region_name, content, design, props, deck_theme):
    """Circular loop of nodes arranged in a ring pattern."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    secondary = _get_secondary(design, deck_theme)

    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})
    nodes = visual_data.get("cycle_nodes", [])

    if not nodes:
        for b in content.get("bullets", [])[:6]:
            nodes.append({"label": b})

    if not nodes:
        return

    import math

    nodes = nodes[:6]
    num = len(nodes)
    accent_rgb = _hex_to_rgb(accent)
    secondary_rgb = _hex_to_rgb(secondary)

    cx = left.inches + width.inches / 2
    cy = top.inches + height.inches / 2
    radius = min(width.inches, height.inches) * 0.35
    node_w, node_h = 1.5, 0.8

    for i, node in enumerate(nodes):
        angle = 2 * math.pi * i / num - math.pi / 2
        nx = cx + radius * math.cos(angle) - node_w / 2
        ny = cy + radius * math.sin(angle) - node_h / 2

        # Node box (rounded rect)
        shape = slide.shapes.add_shape(
            5, Inches(nx), Inches(ny), Inches(node_w), Inches(node_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb if i % 2 == 0 else secondary_rgb
        shape.line.fill.background()

        label = str(node.get("label", node.get("name", "")))
        _add_textbox(
            slide,
            Inches(nx + 0.05),
            Inches(ny + 0.1),
            Inches(node_w - 0.1),
            Inches(node_h - 0.2),
            label,
            font_size=10,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER,
        )

        # Arrow to next node
        if num > 1:
            next_angle = 2 * math.pi * ((i + 1) % num) / num - math.pi / 2
            mid_angle = (angle + next_angle) / 2
            if next_angle < angle:
                mid_angle += math.pi
            arrow_x = cx + radius * 0.7 * math.cos(mid_angle) - 0.12
            arrow_y = cy + radius * 0.7 * math.sin(mid_angle) - 0.12
            _add_textbox(
                slide,
                Inches(arrow_x),
                Inches(arrow_y),
                Inches(0.3),
                Inches(0.3),
                "↻",
                font_size=14,
                bold=True,
                color=_lighten(accent, 0.4),
                alignment=PP_ALIGN.CENTER,
            )


def _render_comparison_matrix(slide, region_name, content, design, props, deck_theme):
    """Grid of cells with row/column headers — SWOT, peer comparison, etc."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)

    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})
    rows = visual_data.get("row_headers", [])
    cols = visual_data.get("col_headers", [])
    cells = visual_data.get("matrix_cells", [])

    # Fallback: 2×2 from bullets
    if not rows and not cols and content.get("bullets"):
        b = content["bullets"]
        rows = ["Category A", "Category B"]
        cols = ["Aspect 1", "Aspect 2"]
        cells = [
            {"row": 0, "col": 0, "value": b[0] if len(b) > 0 else ""},
            {"row": 0, "col": 1, "value": b[1] if len(b) > 1 else ""},
            {"row": 1, "col": 0, "value": b[2] if len(b) > 2 else ""},
            {"row": 1, "col": 1, "value": b[3] if len(b) > 3 else ""},
        ]

    if not rows or not cols:
        return

    accent_rgb = _hex_to_rgb(accent)
    n_rows = len(rows)
    n_cols = len(cols)
    header_w = 1.5
    cell_w = min(2.5, (width.inches - header_w) / max(n_cols, 1))
    cell_h = min(1.0, (height.inches - 0.5) / max(n_rows + 1, 1))

    # Column headers
    for j, col in enumerate(cols[:4]):
        cx = left.inches + header_w + j * cell_w
        shape = slide.shapes.add_shape(
            1, Inches(cx), top, Inches(cell_w), Inches(cell_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()
        _add_textbox(
            slide,
            Inches(cx + 0.05),
            top + Inches(0.1),
            Inches(cell_w - 0.1),
            Inches(cell_h - 0.2),
            str(col),
            font_size=11,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER,
        )

    # Row headers + cells
    for i, row in enumerate(rows[:4]):
        ry = top.inches + cell_h * (i + 1)
        # Row header
        shape = slide.shapes.add_shape(
            1, left, Inches(ry), Inches(header_w), Inches(cell_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _lighten(accent, 0.7)
        shape.line.fill.background()
        _add_textbox(
            slide,
            left + Inches(0.05),
            Inches(ry + 0.05),
            Inches(header_w - 0.1),
            Inches(cell_h - 0.1),
            str(row),
            font_size=10,
            bold=True,
            color=accent_rgb,
            alignment=PP_ALIGN.CENTER,
        )

        # Data cells
        for j in range(min(n_cols, 4)):
            cx = left.inches + header_w + j * cell_w
            cell_val = ""
            for c in cells:
                if c.get("row") == i and c.get("col") == j:
                    cell_val = str(c.get("value", ""))
                    break
            bg = _lighten(accent, 0.92) if (i + j) % 2 == 0 else _lighten(accent, 0.96)
            shape = slide.shapes.add_shape(
                1, Inches(cx), Inches(ry), Inches(cell_w), Inches(cell_h)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
            shape.line.color.rgb = _lighten(accent, 0.7)
            shape.line.width = Pt(0.5)
            _add_textbox(
                slide,
                Inches(cx + 0.05),
                Inches(ry + 0.05),
                Inches(cell_w - 0.1),
                Inches(cell_h - 0.1),
                cell_val,
                font_size=9,
                bold=False,
                color=RGBColor(0x33, 0x33, 0x33),
                alignment=PP_ALIGN.LEFT,
            )


def _render_icon_fact_grid(slide, region_name, content, design, props, deck_theme):
    """Grid of icon + short fact pairs, 2 or 3 columns."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    icon_hints = design.get("icon_suggestions", [])
    accent_hex = accent.lstrip("#")
    bullets = content.get("bullets", [])[:9]

    if not bullets:
        return

    cols = props.get("columns", 3 if len(bullets) > 4 else 2)
    rows_needed = (len(bullets) + cols - 1) // cols
    cell_w = width.inches / cols
    cell_h = min(1.2, height.inches / max(rows_needed, 1))
    icon_sz = Inches(0.3)

    for i, fact in enumerate(bullets):
        r, c = divmod(i, cols)
        fx = left.inches + c * cell_w
        fy = top.inches + r * cell_h

        icon_path = resolve_icon_for_bullet(icon_hints, i, accent_hex)
        if icon_path and os.path.isfile(icon_path):
            slide.shapes.add_picture(
                icon_path, Inches(fx), Inches(fy + 0.1), icon_sz, icon_sz
            )
        else:
            _add_textbox(
                slide,
                Inches(fx),
                Inches(fy + 0.05),
                icon_sz,
                icon_sz,
                "🔹",
                font_size=14,
                alignment=PP_ALIGN.CENTER,
            )

        _add_textbox(
            slide,
            Inches(fx + 0.4),
            Inches(fy + 0.05),
            Inches(cell_w - 0.5),
            Inches(cell_h - 0.1),
            fact,
            font_size=11,
            bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
        )


def _render_card_grid(slide, region_name, content, design, props, deck_theme):
    """Collection of styled cards — products, portfolio items, offerings."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    secondary = _get_secondary(design, deck_theme)

    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})
    card_items = visual_data.get("card_items", [])

    if not card_items:
        for b in content.get("bullets", [])[:6]:
            card_items.append({"title": b, "body": ""})

    if not card_items:
        return

    cards = card_items[:6]
    num = len(cards)
    cols = min(3, num)
    rows = (num + cols - 1) // cols
    card_w = min(3.5, (width.inches - 0.2 * (cols - 1)) / cols)
    card_h = min(2.2, (height.inches - 0.2 * (rows - 1)) / rows)
    total_w = cols * card_w + (cols - 1) * 0.2
    start_x = left.inches + (width.inches - total_w) / 2

    accent_rgb = _hex_to_rgb(accent)
    light_bg = _lighten(accent, 0.92)

    for i, card in enumerate(cards):
        r, c = divmod(i, cols)
        cx = start_x + c * (card_w + 0.2)
        cy = top.inches + r * (card_h + 0.2)

        shape = slide.shapes.add_shape(
            5, Inches(cx), Inches(cy), Inches(card_w), Inches(card_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = light_bg
        shape.line.color.rgb = _lighten(accent, 0.6)
        shape.line.width = Pt(1)
        _apply_shadow(shape, blur_pt=5, offset_pt=2, opacity_pct=20)

        # Top accent bar on card (gradient for visual richness)
        bar = slide.shapes.add_shape(
            1, Inches(cx), Inches(cy), Inches(card_w), Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_rgb
        bar.line.fill.background()

        title = str(card.get("title", card.get("label", "")))
        _add_textbox(
            slide,
            Inches(cx + 0.1),
            Inches(cy + 0.15),
            Inches(card_w - 0.2),
            Inches(0.4),
            title,
            font_size=12,
            bold=True,
            color=accent_rgb,
            alignment=PP_ALIGN.LEFT,
        )

        body = str(card.get("body", card.get("description", "")))
        if body:
            _add_textbox(
                slide,
                Inches(cx + 0.1),
                Inches(cy + 0.55),
                Inches(card_w - 0.2),
                Inches(card_h - 0.7),
                body,
                font_size=9,
                bold=False,
                color=RGBColor(0x55, 0x55, 0x55),
            )


def _render_value_chain(slide, region_name, content, design, props, deck_theme):
    """Linear chain of stages connected by chevron arrows."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)

    visual_intent = content.get("visual_intent", {})
    visual_data = visual_intent.get("visual_data", {})
    steps = visual_data.get("process_steps", [])

    if not steps:
        for i, b in enumerate(content.get("bullets", [])[:7]):
            steps.append({"label": b, "order": i + 1})

    if not steps:
        return

    # Reuse process_flow with chevron styling
    _render_process_flow(slide, region_name, content, design, props, deck_theme)


def _render_quote_callout(slide, region_name, content, design, props, deck_theme):
    """Highlighted quote or tagline in a styled card."""
    left, top, width, height = _resolve_region(region_name)
    accent = _get_accent(design, deck_theme)
    accent_rgb = _hex_to_rgb(accent)

    text = props.get("text", "")
    if not text:
        text = content.get("subtitle", "")
    if not text and content.get("bullets"):
        text = content["bullets"][0]

    # Quote card background
    card_w = min(width.inches, 9.0)
    card_h = min(height.inches, 2.5)
    card_x = left.inches + (width.inches - card_w) / 2
    card_y = top.inches + (height.inches - card_h) / 2

    shape = slide.shapes.add_shape(
        5, Inches(card_x), Inches(card_y), Inches(card_w), Inches(card_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _lighten(accent, 0.92)
    shape.line.fill.background()

    # Left accent bar
    bar = slide.shapes.add_shape(
        1, Inches(card_x), Inches(card_y), Inches(0.08), Inches(card_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_rgb
    bar.line.fill.background()

    # Quote mark
    _add_textbox(
        slide,
        Inches(card_x + 0.25),
        Inches(card_y + 0.1),
        Inches(0.5),
        Inches(0.5),
        "\u201c",
        font_size=36,
        bold=True,
        color=_lighten(accent, 0.4),
    )

    # Quote text
    _add_textbox(
        slide,
        Inches(card_x + 0.4),
        Inches(card_y + 0.5),
        Inches(card_w - 0.8),
        Inches(card_h - 0.8),
        text,
        font_size=16,
        bold=False,
        color=RGBColor(0x33, 0x33, 0x33),
        alignment=PP_ALIGN.LEFT,
        font_name=_get_font(design, "body", deck_theme)["font"],
    )


def _render_section_divider(slide, region_name, content, design, props, deck_theme):
    """Branded section transition — gradient background + large title + decorative elements."""
    accent = _get_accent(design, deck_theme)
    secondary = _get_secondary(design, deck_theme)
    accent_rgb = _hex_to_rgb(accent)

    # Full-slide gradient background
    bg_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    _apply_gradient_fill(bg_shape, accent, secondary, angle=5400000)  # top→bottom
    bg_shape.line.fill.background()

    # Decorative translucent circles
    _add_decorative_circles(slide, accent, deck_theme)

    # Thin horizontal rule
    rule = slide.shapes.add_shape(
        1, Inches(3.5), Inches(4.1), Inches(6.3), Inches(0.04)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    rule.line.fill.background()

    title = content.get("title", "")
    _add_textbox(
        slide,
        Inches(1.5),
        Inches(2.3),
        Inches(10.3),
        Inches(1.8),
        title,
        font_size=42,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        alignment=PP_ALIGN.CENTER,
        font_name=_get_font(design, "title", deck_theme)["font"],
    )

    subtitle = content.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide,
            Inches(2.0),
            Inches(4.4),
            Inches(9.3),
            Inches(1.0),
            subtitle,
            font_size=16,
            bold=False,
            color=_lighten(accent, 0.85),
            alignment=PP_ALIGN.CENTER,
        )


def _render_noop(slide, region_name, content, design, props, deck_theme):
    """Placeholder for components that delegate to the main renderer."""
    pass


# ---------------------------------------------------------------------------
# Modern / Gamma-style component renderers
# ---------------------------------------------------------------------------


def _render_full_bleed_image(slide, region_name, content, design, props, deck_theme):
    """Full-slide image background (image placed by pptx_renderer) + dark overlay + text.

    The actual image placement is handled by pptx_renderer when it detects
    background_treatment == 'full_bleed_image'. This renderer adds the
    semi-transparent darkening overlay and text elements on top.
    """
    accent = _get_accent(design, deck_theme)

    # Dark overlay for text legibility (rendered as a low-opacity dark rectangle)
    overlay = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    overlay.line.fill.background()
    # Set 50% opacity via XML
    spPr = overlay._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            etree.SubElement(srgb, qn("a:alpha"), attrib={"val": "50000"})

    # Hero text overlay
    text = (
        props.get("text")
        or content.get("subtitle", "")
        or content.get("key_takeaway", "")
    )
    if not text and content.get("bullets"):
        text = content["bullets"][0]
    if text:
        title_font = _get_font(design, "title", deck_theme)
        _add_textbox(
            slide,
            Inches(1.5),
            Inches(2.5),
            Inches(10.3),
            Inches(2.5),
            text,
            font_size=36,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER,
            font_name=title_font["font"],
            anchor=MSO_ANCHOR.MIDDLE,
        )


def _render_split_hero(slide, region_name, content, design, props, deck_theme):
    """Asymmetric image + text split. Image is placed by pptx_renderer;
    this renderer handles the text content side with proper formatting.
    """
    accent = _get_accent(design, deck_theme)
    body_spec = _get_font(design, "body", deck_theme)
    title_font = _get_font(design, "title", deck_theme)
    # Determine which side has content (opposite of image)
    side = props.get("image_side", "left")  # where image goes
    if side == "left":
        text_region = "content_right"
    else:
        text_region = "content_left"
    left, top, width, height = _resolve_region(text_region)

    # Render title in the text region
    title = content.get("title", "")
    if title:
        _add_textbox(
            slide,
            left,
            top,
            width,
            Inches(1.0),
            title,
            font_size=title_font["size"],
            bold=True,
            color=_hex_to_rgb(accent),
            font_name=title_font["font"],
        )

    # Render bullets or subtitle below
    bullets = content.get("bullets", [])
    if bullets:
        y_offset = top.inches + 1.3
        for i, bullet in enumerate(bullets[:6]):
            _add_textbox(
                slide,
                Inches(left.inches + 0.1),
                Inches(y_offset + i * 0.44),
                Inches(width.inches - 0.2),
                Inches(0.4),
                f"• {bullet}",
                font_size=body_spec["size"],
                color=_hex_to_rgb(body_spec["color"]),
                font_name=body_spec["font"],
            )
    elif content.get("subtitle"):
        _add_textbox(
            slide,
            left,
            Inches(top.inches + 1.3),
            width,
            Inches(2.0),
            content["subtitle"],
            font_size=body_spec["size"],
            color=_hex_to_rgb(body_spec["color"]),
            font_name=body_spec["font"],
        )


def _render_stat_wall(slide, region_name, content, design, props, deck_theme):
    """Oversized stat number + supporting context below.

    Renders a single dominant metric in huge typography, plus optional
    supporting text underneath for context.
    """
    accent = _get_accent(design, deck_theme)
    accent_rgb = _hex_to_rgb(accent)
    title_font = _get_font(design, "title", deck_theme)
    body_spec = _get_font(design, "body", deck_theme)

    # Extract the stat value from props, supporting_data, or title
    stat_value = props.get("stat_value", "")
    stat_label = props.get("stat_label", "")
    if not stat_value:
        supporting = content.get("supporting_data", [])
        if supporting:
            parts = supporting[0].strip().split(" ", 1)
            stat_value = parts[0]
            stat_label = parts[1] if len(parts) > 1 else ""
    if not stat_value:
        stat_value = content.get("title", "")

    # Big stat — centered, oversized
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(12.3),
        Inches(2.5),
        stat_value,
        font_size=72,
        bold=True,
        color=accent_rgb,
        alignment=PP_ALIGN.CENTER,
        font_name=title_font["font"],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Stat label below
    if stat_label:
        _add_textbox(
            slide,
            Inches(2.0),
            Inches(4.0),
            Inches(9.3),
            Inches(0.8),
            stat_label,
            font_size=22,
            bold=False,
            color=RGBColor(0x55, 0x55, 0x55),
            alignment=PP_ALIGN.CENTER,
            font_name=body_spec["font"],
        )

    # Supporting context / bullets
    bullets = content.get("bullets", [])
    if bullets:
        context_text = " · ".join(bullets[:3])
        _add_textbox(
            slide,
            Inches(1.5),
            Inches(5.2),
            Inches(10.3),
            Inches(1.5),
            context_text,
            font_size=14,
            color=RGBColor(0x66, 0x66, 0x66),
            alignment=PP_ALIGN.CENTER,
            font_name=body_spec["font"],
        )


def _render_pull_quote(slide, region_name, content, design, props, deck_theme):
    """Editorial-style large-font quote with attribution.

    Bigger and bolder than quote_callout — designed for bold mood slides.
    """
    accent = _get_accent(design, deck_theme)
    accent_rgb = _hex_to_rgb(accent)
    title_font = _get_font(design, "title", deck_theme)

    text = props.get("text") or content.get("subtitle", "")
    if not text and content.get("bullets"):
        text = content["bullets"][0]
    attribution = props.get("attribution", "")

    # Large opening quote mark
    _add_textbox(
        slide,
        Inches(1.0),
        Inches(1.5),
        Inches(1.0),
        Inches(1.0),
        "\u201c",
        font_size=80,
        bold=True,
        color=_lighten(accent, 0.5),
        alignment=PP_ALIGN.LEFT,
        font_name=title_font["font"],
    )

    # Quote text — large, centered
    _add_textbox(
        slide,
        Inches(1.8),
        Inches(2.5),
        Inches(9.7),
        Inches(2.5),
        text,
        font_size=28,
        bold=False,
        color=RGBColor(0x1A, 0x1A, 0x1A),
        alignment=PP_ALIGN.LEFT,
        font_name=title_font["font"],
    )

    # Attribution
    if attribution:
        _add_textbox(
            slide,
            Inches(1.8),
            Inches(5.3),
            Inches(9.7),
            Inches(0.5),
            f"— {attribution}",
            font_size=14,
            bold=False,
            color=accent_rgb,
            alignment=PP_ALIGN.LEFT,
        )

    # Decorative bottom accent bar
    bar = slide.shapes.add_shape(1, Inches(1.8), Inches(6.2), Inches(3.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_rgb
    bar.line.fill.background()


def _render_media_overlay(slide, region_name, content, design, props, deck_theme):
    """Image background with semi-transparent text card overlay.

    Image is placed by pptx_renderer. This renderer adds a frosted-glass
    style card on top with content.
    """
    accent = _get_accent(design, deck_theme)
    body_spec = _get_font(design, "body", deck_theme)
    title_font = _get_font(design, "title", deck_theme)

    # Semi-transparent card overlay (lower-center of slide)
    card_x, card_y = 1.5, 3.8
    card_w, card_h = 10.3, 3.2
    card = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(card_x),
        Inches(card_y),
        Inches(card_w),
        Inches(card_h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.fill.background()
    # Set 85% opacity on the card for a frosted-glass feel
    spPr = card._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            etree.SubElement(srgb, qn("a:alpha"), attrib={"val": "85000"})
    _apply_shadow(card, blur_pt=8, offset_pt=2, opacity_pct=20)

    # Title inside card
    title = content.get("title", "")
    if title:
        _add_textbox(
            slide,
            Inches(card_x + 0.4),
            Inches(card_y + 0.3),
            Inches(card_w - 0.8),
            Inches(0.8),
            title,
            font_size=title_font["size"],
            bold=True,
            color=_hex_to_rgb(accent),
            font_name=title_font["font"],
        )

    # Bullets or subtitle inside card
    bullets = content.get("bullets", [])
    if bullets:
        y_off = card_y + 1.2
        for i, bullet in enumerate(bullets[:4]):
            _add_textbox(
                slide,
                Inches(card_x + 0.5),
                Inches(y_off + i * 0.42),
                Inches(card_w - 1.0),
                Inches(0.38),
                f"• {bullet}",
                font_size=body_spec["size"],
                color=RGBColor(0x33, 0x33, 0x33),
                font_name=body_spec["font"],
            )
    elif content.get("subtitle"):
        _add_textbox(
            slide,
            Inches(card_x + 0.4),
            Inches(card_y + 1.2),
            Inches(card_w - 0.8),
            Inches(1.5),
            content["subtitle"],
            font_size=body_spec["size"],
            color=RGBColor(0x33, 0x33, 0x33),
            font_name=body_spec["font"],
        )


# ---------------------------------------------------------------------------
# Component dispatcher
# ---------------------------------------------------------------------------

COMPONENT_RENDERERS = {
    "hero_text": _render_hero_text,
    "bullet_list": _render_bullet_list,
    "text_block": _render_text_block,
    "quote_callout": _render_quote_callout,
    "kpi_strip": _render_kpi_strip,
    "chart_panel": _render_noop,  # handled by pptx_renderer's chart logic
    "table_panel": _render_noop,  # handled by pptx_renderer's table logic
    "timeline": _render_timeline,
    "process_flow": _render_process_flow,
    "cycle_loop": _render_cycle_loop,
    "comparison_matrix": _render_comparison_matrix,
    "hierarchy": _render_bullet_list,  # fallback until hierarchy renderer is built
    "network_map": _render_icon_fact_grid,  # fallback
    "value_chain": _render_value_chain,
    "icon_fact_grid": _render_icon_fact_grid,
    "image_panel": _render_noop,  # handled by pptx_renderer's image logic
    "card_grid": _render_card_grid,
    "section_divider": _render_section_divider,
    # Modern / Gamma-style components
    "full_bleed_image": _render_full_bleed_image,
    "split_hero": _render_split_hero,
    "stat_wall": _render_stat_wall,
    "pull_quote": _render_pull_quote,
    "media_overlay": _render_media_overlay,
}


def render_composition(
    slide,
    composition: dict,
    content: dict,
    design: dict,
    deck_theme: dict | None = None,
) -> bool:
    """Render all components in a composition onto a slide.

    Computes density-aware adaptive regions (Phase 4B) before dispatching
    each component.  Returns True if at least one component was rendered.
    """
    global _active_region_map

    components = composition.get("components", [])
    if not components:
        return False

    # Phase 4B: compute adaptive regions for this slide
    _active_region_map = compute_adaptive_regions(components, content)

    rendered = 0
    try:
        for comp in components:
            comp_type = comp.get("type", "")
            region_name = comp.get("region", "full")
            props = comp.get("props", {})

            renderer = COMPONENT_RENDERERS.get(comp_type)
            if renderer:
                renderer(slide, region_name, content, design, props, deck_theme)
                rendered += 1
    finally:
        _active_region_map = None

    return rendered > 0
