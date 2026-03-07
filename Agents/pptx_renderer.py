"""
PPTX Renderer — Deterministic slide builder.

Takes slide contents + design specs and produces a .pptx file
using python-pptx. All creative decisions were made by the LLM agents;
this module only executes them.

Visual features:
  - Icon bullets with Unicode symbols
  - Metric cards for KPI slides
  - Accent shapes & decorative elements
  - AI-generated image support
  - Slide footers with company name & page number
  - Enhanced layouts with better spacing
"""

import os
import re
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

from Agents.state import PresentationState
from Agents.icon_manager import resolve_icon_for_bullet


# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Default colors
DEFAULT_ACCENT = "#2C3E50"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MEDIUM = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
NEAR_WHITE = RGBColor(0xF5, 0xF5, 0xF5)

# Emoji icon mapping — maps design-agent icon_suggestions to contextual emoji
ICON_EMOJI_MAP = {
    # Financial / Growth
    "growth": "📈",
    "trending_up": "📈",
    "revenue": "💰",
    "profit": "💰",
    "financial": "📊",
    "chart": "📊",
    "performance": "📊",
    "money": "💰",
    # Corporate
    "building": "🏢",
    "company": "🏢",
    "corporate": "🏢",
    "office": "🏢",
    "team": "👥",
    "people": "👥",
    "management": "👥",
    "leadership": "👥",
    # Product / Industry
    "factory": "🏭",
    "manufacturing": "🏭",
    "product": "📦",
    "supply": "📦",
    "technology": "💡",
    "innovation": "💡",
    "idea": "💡",
    "research": "🔬",
    # Risk / Governance
    "risk": "⚠️",
    "warning": "⚠️",
    "shield": "🛡️",
    "security": "🛡️",
    "governance": "⚖️",
    "compliance": "📋",
    "regulation": "📋",
    # Positive / Achievement
    "check": "✅",
    "success": "✅",
    "strength": "💪",
    "advantage": "💪",
    "star": "⭐",
    "highlight": "⭐",
    "key": "🔑",
    "award": "🏆",
    "target": "🎯",
    "goal": "🎯",
    "strategy": "🎯",
    "focus": "🎯",
    # Market / Global
    "market": "🌐",
    "global": "🌐",
    "world": "🌍",
    "expansion": "🚀",
    "competitor": "⚔️",
    "peer": "📊",
    "benchmark": "📊",
    # General
    "arrow": "➡️",
    "next": "➡️",
    "info": "ℹ️",
    "note": "📝",
    "timeline": "⏳",
    "history": "📅",
    "milestone": "🏁",
    "location": "📍",
    "client": "🤝",
    "partner": "🤝",
}

# Default bullet emoji cycle (used when no icon_suggestions available)
DEFAULT_BULLET_EMOJI = ["🔹", "🔸", "🔹", "🔸", "🔹", "🔸"]


# Default text hierarchy fallback (used when design spec lacks entries)
_DEFAULT_HIERARCHY = {
    "title": {"size": 28, "bold": True, "font": "Calibri", "color": "#1a1a1a"},
    "subtitle": {"size": 16, "bold": False, "font": "Calibri", "color": "#555555"},
    "body": {"size": 14, "bold": False, "font": "Calibri", "color": "#333333"},
    "caption": {"size": 10, "bold": False, "font": "Calibri", "color": "#888888"},
    "metric": {
        "size": 36,
        "bold": True,
        "font": "Calibri",
        "color": None,
    },  # None → accent
}


def _font(design: dict, element: str) -> dict:
    """Extract font spec for an element from design's text_hierarchy.

    Returns dict with keys: size (int), bold (bool), font (str), color (str|None).
    Falls back to _DEFAULT_HIERARCHY if missing.
    """
    hierarchy = design.get("text_hierarchy", {})
    spec = hierarchy.get(element, {})
    default = _DEFAULT_HIERARCHY.get(element, _DEFAULT_HIERARCHY["body"])
    return {
        "size": spec.get("size", default["size"]),
        "bold": spec.get("bold", default["bold"]),
        "font": spec.get("font", default["font"]),
        "color": spec.get("color", default["color"]),
    }


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        hex_color = "2C3E50"  # fallback
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _lighten(hex_color: str, factor: float = 0.85) -> RGBColor:
    """Return a lightened version of a hex color (for card backgrounds)."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        hex_color = "2C3E50"
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return RGBColor(min(r, 255), min(g, 255), min(b, 255))


def _add_background(slide, hex_color: str):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    bold=False,
    color=DARK,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_footer(
    slide, company_name: str, slide_num: int, total_slides: int, accent: str
):
    """Add a footer bar with company name and page number."""
    # Thin accent bar at bottom
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.15), SLIDE_WIDTH, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(accent)
    bar.line.fill.background()

    # Company name on left
    _add_textbox(
        slide,
        Inches(0.4),
        Inches(7.15),
        Inches(6),
        Inches(0.35),
        company_name,
        font_size=9,
        bold=False,
        color=WHITE,
        alignment=PP_ALIGN.LEFT,
    )

    # Page number on right
    _add_textbox(
        slide,
        Inches(10),
        Inches(7.15),
        Inches(3),
        Inches(0.35),
        f"{slide_num} / {total_slides}",
        font_size=9,
        bold=False,
        color=WHITE,
        alignment=PP_ALIGN.RIGHT,
    )


def _add_accent_bar_top(slide, accent: str):
    """Add a decorative accent line at the top of a slide."""
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _hex_to_rgb(accent)
    line.line.fill.background()


def _add_side_accent(slide, accent: str):
    """Add a thin decorative accent strip on the left edge."""
    strip = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT)
    strip.fill.solid()
    strip.fill.fore_color.rgb = _hex_to_rgb(accent)
    strip.line.fill.background()


def _add_corner_triangle(slide, accent: str):
    """Add a small decorative triangle in the bottom-right corner."""
    # Use a right-triangle freeform approximation via a rectangle rotated
    # python-pptx doesn't support freeform easily, use a small accent rectangle
    rect = slide.shapes.add_shape(
        1, Inches(12.5), Inches(6.5), Inches(0.83), Inches(0.65)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = _lighten(accent, 0.6)
    rect.line.fill.background()
    rect.rotation = 0


def _add_image_to_slide(slide, image_path: str, left, top, width, height):
    """Add an image to the slide if the file exists."""
    if image_path and os.path.isfile(image_path):
        slide.shapes.add_picture(image_path, left, top, width, height)
        return True
    return False


def _resolve_emoji(icon_hints: list, index: int) -> str:
    """Emoji fallback when Lucide icon PNGs are unavailable."""
    if icon_hints:
        for hint in icon_hints:
            key = hint.lower().strip()
            if key in ICON_EMOJI_MAP:
                return ICON_EMOJI_MAP[key]
        for hint in icon_hints:
            key = hint.lower().strip()
            for map_key, emoji in ICON_EMOJI_MAP.items():
                if key in map_key or map_key in key:
                    return emoji
    return DEFAULT_BULLET_EMOJI[index % len(DEFAULT_BULLET_EMOJI)]


def _render_icon_bullets(
    slide,
    bullets: list,
    top_offset,
    accent: str,
    left=Inches(0.7),
    width=Inches(11.5),
    icon_hints: list = None,
    body_spec: dict = None,
):
    """Render bullets with Lucide icon images (or emoji fallback).

    Each bullet is an individual icon image + text box pair, giving
    per-bullet icon resolution from the Lucide library.
    body_spec: font spec dict with keys size, bold, font, color.
    """
    if not bullets:
        return

    bs = body_spec or _DEFAULT_HIERARCHY["body"]
    font_size = bs.get("size", 14)
    font_name = bs.get("font", "Calibri")
    font_color = _hex_to_rgb(bs["color"]) if bs.get("color") else BODY_COLOR

    accent_hex = accent.lstrip("#") if isinstance(accent, str) else "2C3E50"
    icon_size = Inches(0.28)
    icon_gap = Inches(0.12)
    text_left = Inches(left.inches + icon_size.inches + icon_gap.inches)
    text_width = Inches(width.inches - icon_size.inches - icon_gap.inches)
    row_height = Inches(0.44)

    for i, bullet in enumerate(bullets):
        y = Inches(top_offset.inches + i * row_height.inches)

        # Try Lucide PNG icon first
        icon_path = resolve_icon_for_bullet(icon_hints, i, accent_hex)

        if icon_path and os.path.isfile(icon_path):
            slide.shapes.add_picture(icon_path, left, y, icon_size, icon_size)
        else:
            # Emoji text fallback
            emoji = _resolve_emoji(icon_hints, i)
            _add_textbox(
                slide,
                left,
                y,
                icon_size,
                Inches(0.3),
                emoji,
                font_size=13,
                alignment=PP_ALIGN.CENTER,
            )

        # Bullet text
        txBox = slide.shapes.add_textbox(text_left, y, text_width, Inches(0.36))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name


def _render_metric_cards(
    slide, supporting_data: list, accent: str, top=Inches(1.8), metric_spec: dict = None
):
    """Render key metrics as visually prominent cards."""
    if not supporting_data:
        return

    ms = metric_spec or _DEFAULT_HIERARCHY["metric"]
    metric_font = ms.get("font", "Calibri")
    metric_size = ms.get("size", 36)

    # Limit to 4 metric cards max
    metrics = supporting_data[:4]
    num = len(metrics)
    if num == 0:
        return

    card_width = Inches(2.6)
    card_height = Inches(1.5)
    gap = Inches(0.3)
    total_width = num * card_width.inches + (num - 1) * gap.inches
    start_left = (13.333 - total_width) / 2  # center cards

    accent_rgb = _hex_to_rgb(accent)
    metric_color = _hex_to_rgb(ms["color"]) if ms.get("color") else accent_rgb
    light_bg = _lighten(accent, 0.88)

    for i, metric in enumerate(metrics):
        left = Inches(start_left + i * (card_width.inches + gap.inches))

        # Card background
        card = slide.shapes.add_shape(
            5, left, top, card_width, card_height  # rounded rectangle
        )
        card.fill.solid()
        card.fill.fore_color.rgb = light_bg
        card.line.color.rgb = _lighten(accent, 0.6)
        card.line.width = Pt(1)

        # Metric value (the whole string — e.g. "₹2,366 Cr revenue (FY25)")
        # Try to split into number part and label
        parts = metric.strip().split(" ", 1)
        value_text = parts[0] if parts else metric
        label_text = parts[1] if len(parts) > 1 else ""

        # Big number
        _add_textbox(
            slide,
            left + Inches(0.15),
            top + Inches(0.2),
            card_width - Inches(0.3),
            Inches(0.7),
            value_text,
            font_size=min(metric_size, 42),
            bold=True,
            color=metric_color,
            alignment=PP_ALIGN.CENTER,
            font_name=metric_font,
        )

        # Label below
        if label_text:
            _add_textbox(
                slide,
                left + Inches(0.15),
                top + Inches(0.9),
                card_width - Inches(0.3),
                Inches(0.5),
                label_text,
                font_size=11,
                bold=False,
                color=MEDIUM,
                alignment=PP_ALIGN.CENTER,
            )


def _render_title_slide(prs, content: dict, design: dict, image_map: dict = None):
    """Render a full-screen title/cover slide with decorative elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    accent = design.get("color_accent", DEFAULT_ACCENT)
    tf_spec = _font(design, "title")
    sf_spec = _font(design, "subtitle")
    _add_background(slide, accent)

    # Decorative lighter rectangle overlay (top-right area)
    overlay = slide.shapes.add_shape(1, Inches(7), Inches(0), Inches(6.333), Inches(3))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = _lighten(accent, 0.15)
    overlay.line.fill.background()

    # Small decorative square
    sq = slide.shapes.add_shape(1, Inches(11.5), Inches(5.5), Inches(0.6), Inches(0.6))
    sq.fill.solid()
    sq.fill.fore_color.rgb = _lighten(accent, 0.25)
    sq.line.fill.background()

    # Company name / title — use agent-specified font, boost size for cover
    title = content.get("title", "")
    cover_title_size = max(tf_spec["size"], 40)  # cover titles should be big
    _add_textbox(
        slide,
        Inches(1),
        Inches(2.2),
        Inches(11),
        Inches(1.5),
        title,
        font_size=cover_title_size,
        bold=tf_spec["bold"],
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
        font_name=tf_spec["font"],
    )

    # Thin horizontal rule under title
    rule = slide.shapes.add_shape(1, Inches(4), Inches(3.8), Inches(5.3), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = WHITE
    rule.line.fill.background()

    # Subtitle
    subtitle = content.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide,
            Inches(1.5),
            Inches(4.2),
            Inches(10),
            Inches(1),
            subtitle,
            font_size=max(sf_spec["size"], 18),
            bold=sf_spec["bold"],
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            font_name=sf_spec["font"],
        )

    # Key takeaway at bottom
    takeaway = content.get("key_takeaway", "")
    if takeaway and content.get("slide_id") != "cover":
        _add_textbox(
            slide,
            Inches(2),
            Inches(5.5),
            Inches(9),
            Inches(0.8),
            takeaway,
            font_size=14,
            bold=False,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            font_name=sf_spec["font"],
        )

    # Add speaker notes
    notes = content.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _render_section_header(prs, content: dict, design: dict, image_map: dict = None):
    """Render a section header/divider slide with accent bar and decorative elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    accent = design.get("color_accent", DEFAULT_ACCENT)
    tf_spec = _font(design, "title")
    sf_spec = _font(design, "subtitle")

    # Wide accent bar on left
    left_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.5), SLIDE_HEIGHT
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = _hex_to_rgb(accent)
    left_bar.line.fill.background()

    # Light accent block in background
    bg_block = slide.shapes.add_shape(
        1, Inches(0.5), Inches(2), Inches(12.833), Inches(3.5)
    )
    bg_block.fill.solid()
    bg_block.fill.fore_color.rgb = _lighten(accent, 0.92)
    bg_block.line.fill.background()

    # Title — use agent font, boosted for section headers
    title_color = _hex_to_rgb(tf_spec["color"]) if tf_spec.get("color") else DARK
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(2.5),
        Inches(10),
        Inches(1.5),
        content.get("title", ""),
        font_size=max(tf_spec["size"], 36),
        bold=tf_spec["bold"],
        color=title_color,
        font_name=tf_spec["font"],
    )

    # Subtitle
    subtitle = content.get("subtitle", "")
    if subtitle:
        sub_color = _hex_to_rgb(sf_spec["color"]) if sf_spec.get("color") else MEDIUM
        _add_textbox(
            slide,
            Inches(1.2),
            Inches(4.2),
            Inches(10),
            Inches(1),
            subtitle,
            font_size=sf_spec["size"],
            bold=sf_spec["bold"],
            color=sub_color,
            font_name=sf_spec["font"],
        )

    notes = content.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _render_title_content(prs, content: dict, design: dict, image_map: dict = None):
    """Render a standard title + bullet content slide with icon bullets and optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    accent = design.get("color_accent", DEFAULT_ACCENT)
    slide_id = content.get("slide_id", "")
    content_type = content.get("content_type", "bullets")
    icon_hints = design.get("icon_suggestions", [])
    tf_spec = _font(design, "title")
    sf_spec = _font(design, "subtitle")
    bf_spec = _font(design, "body")
    mf_spec = _font(design, "metric")

    # Decorative elements
    _add_side_accent(slide, accent)
    _add_accent_bar_top(slide, accent)

    # Title
    title_color = _hex_to_rgb(tf_spec["color"]) if tf_spec.get("color") else DARK
    _add_textbox(
        slide,
        Inches(0.7),
        Inches(0.6),
        Inches(12),
        Inches(0.8),
        content.get("title", ""),
        font_size=tf_spec["size"],
        bold=tf_spec["bold"],
        color=title_color,
        font_name=tf_spec["font"],
    )

    # Subtitle (if any)
    subtitle = content.get("subtitle", "")
    top_offset = Inches(1.5)
    if subtitle:
        sub_color = _hex_to_rgb(sf_spec["color"]) if sf_spec.get("color") else MEDIUM
        _add_textbox(
            slide,
            Inches(0.7),
            Inches(1.4),
            Inches(12),
            Inches(0.5),
            subtitle,
            font_size=sf_spec["size"],
            bold=sf_spec["bold"],
            color=sub_color,
            font_name=sf_spec["font"],
        )
        top_offset = Inches(2.1)

    # Check for metric cards (if content_type is "metrics" and supporting_data exists)
    supporting_data = content.get("supporting_data", [])
    if content_type == "metrics" and supporting_data:
        _render_metric_cards(
            slide, supporting_data, accent, top=top_offset, metric_spec=mf_spec
        )
        # Put bullets below metric cards
        bullet_top = Inches(top_offset.inches + 1.8)
    else:
        bullet_top = top_offset

    # Check for AI-generated image
    image_path = (image_map or {}).get(slide_id)
    has_image = image_path and os.path.isfile(image_path)

    if has_image:
        # Image on right, bullets on left
        _render_icon_bullets(
            slide,
            content.get("bullets", []),
            bullet_top,
            accent,
            left=Inches(0.7),
            width=Inches(6.5),
            icon_hints=icon_hints,
            body_spec=bf_spec,
        )
        _add_image_to_slide(
            slide,
            image_path,
            Inches(7.8),
            bullet_top,
            Inches(5),
            Inches(3.5),
        )
    else:
        # Full-width bullets
        _render_icon_bullets(
            slide,
            content.get("bullets", []),
            bullet_top,
            accent,
            icon_hints=icon_hints,
            body_spec=bf_spec,
        )

    # Key takeaway at bottom
    cf_spec = _font(design, "caption")
    takeaway = content.get("key_takeaway", "")
    if takeaway:
        # Takeaway card
        card = slide.shapes.add_shape(
            5, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6)  # rounded rect
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _lighten(accent, 0.88)
        card.line.fill.background()

        _add_textbox(
            slide,
            Inches(0.7),
            Inches(6.25),
            Inches(11.9),
            Inches(0.5),
            f"⚡ KEY INSIGHT:  {takeaway}",
            font_size=12,
            bold=True,
            color=_hex_to_rgb(accent),
            font_name=bf_spec["font"],
        )

    notes = content.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _render_two_column(prs, content: dict, design: dict, image_map: dict = None):
    """Render an enhanced two-column layout slide with decorative divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    accent = design.get("color_accent", DEFAULT_ACCENT)
    icon_hints = design.get("icon_suggestions", [])
    tf_spec = _font(design, "title")
    sf_spec = _font(design, "subtitle")
    bf_spec = _font(design, "body")

    # Decorative elements
    _add_side_accent(slide, accent)
    _add_accent_bar_top(slide, accent)

    # Title
    title_color = _hex_to_rgb(tf_spec["color"]) if tf_spec.get("color") else DARK
    _add_textbox(
        slide,
        Inches(0.7),
        Inches(0.6),
        Inches(12),
        Inches(0.8),
        content.get("title", ""),
        font_size=tf_spec["size"],
        bold=tf_spec["bold"],
        color=title_color,
        font_name=tf_spec["font"],
    )

    # Split bullets into two columns
    bullets = content.get("bullets", [])
    mid = (len(bullets) + 1) // 2
    left_bullets = bullets[:mid]
    right_bullets = bullets[mid:]

    # Left column with icon bullets
    _render_icon_bullets(
        slide,
        left_bullets,
        Inches(1.8),
        accent,
        left=Inches(0.7),
        width=Inches(5.5),
        icon_hints=icon_hints,
        body_spec=bf_spec,
    )

    # Right column with icon bullets
    _render_icon_bullets(
        slide,
        right_bullets,
        Inches(1.8),
        accent,
        left=Inches(7.0),
        width=Inches(5.5),
        icon_hints=icon_hints,
        body_spec=bf_spec,
    )

    # Vertical divider (styled)
    divider = slide.shapes.add_shape(
        1, Inches(6.45), Inches(1.8), Inches(0.04), Inches(4.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = _hex_to_rgb(accent)
    divider.line.fill.background()

    # Optional column headers from subtitle
    subtitle = content.get("subtitle", "")
    if subtitle and "|" in subtitle:
        parts = subtitle.split("|", 1)
        sub_color = (
            _hex_to_rgb(sf_spec["color"])
            if sf_spec.get("color")
            else _hex_to_rgb(accent)
        )
        _add_textbox(
            slide,
            Inches(0.7),
            Inches(1.4),
            Inches(5.5),
            Inches(0.4),
            parts[0].strip(),
            font_size=sf_spec["size"],
            bold=True,
            color=sub_color,
            font_name=sf_spec["font"],
        )
        _add_textbox(
            slide,
            Inches(7.0),
            Inches(1.4),
            Inches(5.5),
            Inches(0.4),
            parts[1].strip(),
            font_size=sf_spec["size"],
            bold=True,
            color=sub_color,
            font_name=sf_spec["font"],
        )

    notes = content.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _render_chart_slide(prs, content: dict, design: dict, image_map: dict = None):
    """Render a slide with a chart and decorative elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    accent = design.get("color_accent", DEFAULT_ACCENT)
    tf_spec = _font(design, "title")

    # Decorative elements
    _add_side_accent(slide, accent)
    _add_accent_bar_top(slide, accent)

    # Title
    title_color = _hex_to_rgb(tf_spec["color"]) if tf_spec.get("color") else DARK
    _add_textbox(
        slide,
        Inches(0.7),
        Inches(0.6),
        Inches(12),
        Inches(0.8),
        content.get("title", ""),
        font_size=tf_spec["size"],
        bold=tf_spec["bold"],
        color=title_color,
        font_name=tf_spec["font"],
    )

    chart_data_spec = design.get("chart_data")
    chart_type = design.get("chart_type", "bar")

    if chart_data_spec and chart_type in ("bar", "line", "pie"):
        _add_chart(slide, chart_type, chart_data_spec, accent)
    else:
        # Fallback: render as title_content if no valid chart data
        _render_bullets_on_slide(slide, content)

    # Key takeaway
    takeaway = content.get("key_takeaway", "")
    if takeaway:
        card = slide.shapes.add_shape(
            5, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _lighten(accent, 0.88)
        card.line.fill.background()

        _add_textbox(
            slide,
            Inches(0.7),
            Inches(6.25),
            Inches(11.9),
            Inches(0.5),
            f"⚡ KEY INSIGHT:  {takeaway}",
            font_size=12,
            bold=True,
            color=_hex_to_rgb(accent),
        )

    notes = content.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_chart(slide, chart_type: str, chart_data_spec: dict, accent: str):
    """Add a chart to the slide from structured chart data."""
    categories = chart_data_spec.get("categories", [])
    series_list = chart_data_spec.get("series", [])
    chart_title = chart_data_spec.get("title", "")

    if not categories or not series_list:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series in series_list:
        values = series.get("values", [])
        # Ensure all values are numeric
        clean_values = []
        for v in values:
            try:
                clean_values.append(float(v))
            except (ValueError, TypeError):
                clean_values.append(0)
        chart_data.add_series(series.get("name", "Data"), clean_values)

    # Map chart type
    ct_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }
    xl_chart_type = ct_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_frame = slide.shapes.add_chart(
        xl_chart_type,
        Inches(0.8),
        Inches(1.6),
        Inches(11.5),
        Inches(4.8),
        chart_data,
    )

    chart = chart_frame.chart

    # Always show legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = MEDIUM

    # Extract unit from chart title (text in parentheses, e.g. "Revenue (₹ Cr)")
    unit_match = re.search(r"\(([^)]+)\)", chart_title)
    y_unit = unit_match.group(1) if unit_match else ""

    # Add axis titles for non-pie charts
    if chart_type != "pie":
        # Infer X-axis label from category values
        x_label = ""
        if categories:
            sample = str(categories[0])
            if re.match(r"(FY|CY|Q[1-4]|H[12]|20\d{2}|19\d{2})", sample):
                x_label = "Period"

        # Category axis (X)
        category_axis = chart.category_axis
        if x_label:
            category_axis.has_title = True
            ax_title = category_axis.axis_title.text_frame.paragraphs[0]
            ax_title.text = x_label
            ax_title.font.size = Pt(10)
            ax_title.font.color.rgb = MEDIUM
            ax_title.font.name = "Calibri"
        category_axis.tick_labels.font.size = Pt(9)
        category_axis.tick_labels.font.color.rgb = MEDIUM

        # Value axis (Y)
        value_axis = chart.value_axis
        if y_unit:
            value_axis.has_title = True
            vy_title = value_axis.axis_title.text_frame.paragraphs[0]
            vy_title.text = y_unit
            vy_title.font.size = Pt(10)
            vy_title.font.color.rgb = MEDIUM
            vy_title.font.name = "Calibri"
        value_axis.tick_labels.font.size = Pt(9)
        value_axis.tick_labels.font.color.rgb = MEDIUM
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Style the chart series colors
    accent_rgb = _hex_to_rgb(accent)
    if chart_type != "pie":
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            if i == 0:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = accent_rgb
            else:
                # Secondary series in lighter color
                lighter = RGBColor(
                    min(accent_rgb[0] + 60, 255),
                    min(accent_rgb[1] + 60, 255),
                    min(accent_rgb[2] + 60, 255),
                )
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = lighter


def _render_bullets_on_slide(slide, content: dict):
    """Fallback: render bullets when chart data is unavailable."""
    bullets = content.get("bullets", [])
    if not bullets:
        return

    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.8), Inches(11.5), Inches(4.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(14)
        p.font.color.rgb = BODY_COLOR
        p.font.name = "Calibri"
        p.space_after = Pt(8)


def _render_table_slide(prs, content: dict, design: dict, image_map: dict = None):
    """Render a slide with a data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    accent = design.get("color_accent", DEFAULT_ACCENT)
    tf_spec = _font(design, "title")

    # Decorative elements
    _add_side_accent(slide, accent)
    _add_accent_bar_top(slide, accent)

    # Title
    title_color = _hex_to_rgb(tf_spec["color"]) if tf_spec.get("color") else DARK
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.6),
        Inches(12),
        Inches(0.8),
        content.get("title", ""),
        font_size=tf_spec["size"],
        bold=tf_spec["bold"],
        color=title_color,
        font_name=tf_spec["font"],
    )

    chart_data_spec = design.get("chart_data", {})
    categories = chart_data_spec.get("categories", [])
    series_list = chart_data_spec.get("series", [])

    if categories and series_list:
        # Create table
        rows = len(series_list) + 1  # header + data rows
        cols = len(categories) + 1  # label column + categories
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.0)
        )
        table = table_shape.table

        # Header row
        table.cell(0, 0).text = ""
        for j, cat in enumerate(categories):
            table.cell(0, j + 1).text = str(cat)

        # Data rows
        for i, series in enumerate(series_list):
            table.cell(i + 1, 0).text = series.get("name", "")
            values = series.get("values", [])
            for j, val in enumerate(values):
                if j + 1 < cols:
                    table.cell(i + 1, j + 1).text = str(val)

        # Style header row
        for j in range(cols):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(accent)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = WHITE
                paragraph.font.bold = True
                paragraph.font.size = Pt(11)
    else:
        # Fallback to bullets
        _render_bullets_on_slide(slide, content)

    notes = content.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# Layout dispatcher
LAYOUT_RENDERERS = {
    "title_slide": _render_title_slide,
    "section_header": _render_section_header,
    "title_content": _render_title_content,
    "two_column": _render_two_column,
    "chart_slide": _render_chart_slide,
    "blank": _render_title_content,  # fallback to standard layout
}


def _match_design_spec(slide_id: str, design_specs: list[dict]) -> dict:
    """Find the design spec matching a slide_id."""
    for spec in design_specs:
        if spec.get("slide_id") == slide_id:
            return spec
    # Fallback default
    return {
        "slide_id": slide_id,
        "layout": "title_content",
        "chart_type": "none",
        "chart_data": None,
        "color_accent": DEFAULT_ACCENT,
        "icon_suggestions": [],
        "text_hierarchy": {},
        "visual_balance": "balanced",
    }


def render(state: PresentationState) -> dict:
    """
    LangGraph node function (deterministic — no LLM).
    Renders all slides into a .pptx file.
    """
    company_name = state["company_name"]
    slide_contents = state["slide_contents"]
    design_specs = state["design_specs"]
    image_map = state.get("image_map", {})

    print(f"\n{'='*60}")
    print(f"[Renderer] Building PPTX — {company_name}")
    print(f"{'='*60}")

    # Create presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total_slides = len(slide_contents)

    for idx, content in enumerate(slide_contents):
        slide_id = content.get("slide_id", "unknown")
        design = _match_design_spec(slide_id, design_specs)
        layout = design.get("layout", "title_content")
        chart_type = design.get("chart_type", "none")

        # Handle table type as a special layout
        if chart_type == "table":
            renderer = _render_table_slide
        else:
            renderer = LAYOUT_RENDERERS.get(layout, _render_title_content)

        print(
            f"  Rendering [{slide_id}] → {layout}"
            + (f" + {chart_type} chart" if chart_type and chart_type != "none" else "")
            + (f" + image" if slide_id in image_map else "")
        )
        renderer(prs, content, design, image_map)

        # Add footer to all slides except title_slide
        if layout != "title_slide":
            slide = prs.slides[-1]  # last added slide
            _add_footer(
                slide,
                company_name,
                idx + 1,
                total_slides,
                design.get("color_accent", DEFAULT_ACCENT),
            )

    # Save output
    output_dir = os.path.join("Agents", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    safe_name = company_name.replace(" ", "_").replace("/", "-")
    output_path = os.path.join(output_dir, f"{safe_name}_presentation.pptx")
    prs.save(output_path)
    print(f"\n  Saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return {"pptx_path": output_path}
